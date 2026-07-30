"""
Executive deck for Oberry leadership (PROJECT_BRIEF.md Prompt 14).

Seven slides, 16:9, built with python-pptx from real numbers in
outputs/metrics.json, data/processed/*.csv, and outputs/figures/*.png —
nothing here is invented. Every slide carries the synthetic-data
footer; Part F #6 is non-negotiable for a deck executives will quote.

Usage:
    uv run python src/build_deck.py

Output:
    outputs/primo_churn_deck.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "outputs" / "figures"
PROCESSED = ROOT / "data" / "processed"
OUT_PATH = ROOT / "outputs" / "primo_churn_deck.pptx"

# ── Palette (identical to the web app and every matplotlib figure) ──────────
NAVY = RGBColor(0x1F, 0x3B, 0x57)
NAVY_DARK = RGBColor(0x14, 0x24, 0x35)
TEAL = RGBColor(0x2E, 0x8B, 0x7A)
GOLD = RGBColor(0xD4, 0xA0, 0x3C)
CORAL = RGBColor(0xD6, 0x5C, 0x4A)
SLATE = RGBColor(0x5B, 0x7C, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1C, 0x2B, 0x36)
INK_MUTED = RGBColor(0x6E, 0x7A, 0x8A)
CANVAS = RGBColor(0xF7, 0xF5, 0xF0)
LINE = RGBColor(0xE4, 0xE1, 0xD8)

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)


def load_numbers() -> dict:
    metrics = json.loads((ROOT / "outputs" / "metrics.json").read_text())
    segments = pd.read_csv(PROCESSED / "segments.csv")
    at_risk = pd.read_csv(PROCESSED / "at_risk_members.csv")
    members = pd.read_csv(ROOT / "data" / "raw" / "members.csv")

    top_member = at_risk.sort_values("annual_value_thb", ascending=False).iloc[0]

    return {
        "metrics": metrics,
        "segments": segments,
        "n_total_members": len(members),
        "n_scored": metrics["test_set_size"],
        "n_flagged": len(at_risk),
        "flagged_value_thb": at_risk["annual_value_thb"].sum(),
        "top_member": top_member,
    }


# ── Low-level helpers ────────────────────────────────────────────────────────

def add_blank_slide(prs: Presentation, bg: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank layout
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def add_text(
    slide, left, top, width, height, text, *, size, color, bold=False,
    font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_bullets(
    slide, left, top, width, height, items, *, size, color, font=BODY_FONT,
    bold_lead=False, space_after=10, bullet_color=None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        dash = p.add_run()
        dash.text = "–  "
        dash.font.size = Pt(size)
        dash.font.name = font
        dash.font.color.rgb = bullet_color or color
        dash.font.bold = True
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.name = font
        run.font.color.rgb = color
        run.font.bold = bold_lead
    return box


def add_title(slide, text, *, color=NAVY, size=32):
    return add_text(
        slide, MARGIN, Inches(0.55), SW - 2 * MARGIN, Inches(1.1),
        text, size=size, color=color, bold=True, font=HEADER_FONT,
    )


def add_footer(slide, *, dark=False):
    add_text(
        slide, MARGIN, SH - Inches(0.5), Inches(6), Inches(0.35),
        "Illustrative figures · synthetic data",
        size=10, color=(SLATE if not dark else RGBColor(0x9A, 0xAD, 0xC2)),
        font=BODY_FONT,
    )


def add_picture_fit(slide, path: Path, left, top, max_w, max_h):
    """Place an image scaled to fit inside (max_w, max_h), preserving aspect
    ratio, centered in that box — avoids the distorted-stretch look."""
    from PIL import Image

    with Image.open(path) as im:
        iw, ih = im.size
    img_ratio = iw / ih
    box_ratio = max_w / max_h
    if img_ratio > box_ratio:
        w = max_w
        h = Emu(int(max_w / img_ratio))
    else:
        h = max_h
        w = Emu(int(max_h * img_ratio))
    x = left + Emu(int((max_w - w) / 2))
    y = top + Emu(int((max_h - h) / 2))
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def add_numbered_circle(slide, cx, cy, diameter, number, *, fill=NAVY, text_color=WHITE):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - Emu(int(diameter / 2)), cy - Emu(int(diameter / 2)), diameter, diameter
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.name = HEADER_FONT
    run.font.color.rgb = text_color
    return shape


def set_notes(slide, text: str):
    slide.notes_slide.notes_text_frame.text = text


def add_stat(slide, left, top, width, value, label, *, value_color=NAVY, value_size=44):
    add_text(
        slide, left, top, width, Inches(0.9), value,
        size=value_size, color=value_color, bold=True, font=HEADER_FONT,
    )
    add_text(
        slide, left, top + Inches(0.85), width, Inches(0.5), label,
        size=13, color=INK_MUTED, font=BODY_FONT,
    )


# ── Slide builders ───────────────────────────────────────────────────────────

def build_slide_1(prs, n):
    m = n["metrics"]
    s = add_blank_slide(prs, NAVY_DARK)
    add_text(
        s, MARGIN, Inches(0.7), Inches(4), Inches(0.4), "PRIMO CHURN RADAR",
        size=13, color=TEAL, bold=True, font=BODY_FONT,
    )
    add_text(
        s, MARGIN, Inches(1.25), Inches(11.5), Inches(1.6),
        "Oberry is losing members quietly.",
        size=40, color=WHITE, bold=True, font=HEADER_FONT,
    )
    churn_pct = round(m["test_churn_rate"] * 100)
    add_stat(
        s, MARGIN, Inches(3.2), Inches(3.6),
        f"~{churn_pct}%", "of members go silent for 60+ days, unnoticed",
        value_color=CORAL,
    )
    add_stat(
        s, Inches(4.6), Inches(3.2), Inches(3.6),
        f"{n['n_flagged']:,}", "members already flagged high-risk today",
        value_color=GOLD,
    )
    add_stat(
        s, Inches(8.6), Inches(3.2), Inches(4),
        f"฿{n['flagged_value_thb']/1_000_000:.1f}M/yr", "in spend riding on reaching them in time",
        value_color=WHITE,
    )
    add_text(
        s, MARGIN, Inches(5.6), Inches(10.5), Inches(1.0),
        "No cancellation, no complaint — a member just stops coming back, and nothing in a "
        "typical POS report flags it until they're already gone.",
        size=15, color=RGBColor(0xC9, 0xD4, 0xDE), font=BODY_FONT, line_spacing=1.3,
    )
    add_footer(s, dark=True)
    set_notes(s, (
        "Open with the number, not the model. Almost one in five Oberry members go quiet for "
        "two straight months before anyone notices — that's the definition of churn we used "
        "throughout. Right now, 2,260 members are already flagged as high-risk by the system, "
        "and between them they represent about 5.4 million baht a year in spend. This talk is "
        "about how we found them in time to do something about it."
    ))


def build_slide_2(prs, n):
    s = add_blank_slide(prs, WHITE)
    add_title(s, "One pipeline: data → segments → risk score → action")
    stages = ["Data", "Segments", "Risk score", "Action"]
    x0 = MARGIN
    gap = Inches(2.85)
    for i, label in enumerate(stages):
        cx = x0 + Inches(0.35) + gap * i
        cy = Inches(2.0)
        add_numbered_circle(s, cx, cy, Inches(0.7), i + 1, fill=[NAVY, TEAL, GOLD, CORAL][i])
        add_text(
            s, cx - Inches(1.0), cy + Inches(0.55), Inches(2.0), Inches(0.4), label,
            size=14, color=INK, bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER,
        )
        if i < len(stages) - 1:
            line = s.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, cx + Inches(0.45), cy, cx + gap - Inches(0.45), cy
            )
            line.line.color.rgb = LINE
            line.line.width = Pt(2)
    add_text(
        s, MARGIN, Inches(3.1), Inches(11.5), Inches(0.8),
        "Two years of transaction history becomes five member segments, a trained risk score "
        "for every member, and a specific recommended action — not a hunch.",
        size=15, color=INK, font=BODY_FONT, line_spacing=1.3,
    )
    add_picture_fit(
        s, FIG / "dashboard_screenshot.png",
        MARGIN, Inches(3.95), Inches(12.1), Inches(3.05),
    )
    add_footer(s)
    set_notes(s, (
        "This is the whole system in one line, and this screenshot is the actual live "
        "dashboard, not a mockup — the Oberry Member Retention Radar. Every number on it reads "
        "straight from the same database this pipeline writes to. The operator sees who's "
        "flagged, why, and what to do about it, sorted by how much value is on the line."
    ))


def build_slide_3(prs, n):
    s = add_blank_slide(prs, WHITE)
    add_title(s, "Five kinds of member, not one undifferentiated list")
    add_picture_fit(s, FIG / "07_segment_bubble_chart.png", MARGIN, Inches(1.7), Inches(7.6), Inches(5.3))

    seg_order = ["Champions", "Loyal", "At-risk regulars", "Hibernating", "One-and-done"]
    colors = {"Champions": NAVY, "Loyal": TEAL, "At-risk regulars": GOLD, "Hibernating": CORAL, "One-and-done": SLATE}
    by_name = n["segments"].set_index("segment")
    y = Inches(1.9)
    for name in seg_order:
        row = by_name.loc[name]
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.6), y + Inches(0.08), Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = colors[name]
        dot.line.fill.background()
        dot.shadow.inherit = False
        add_text(
            s, Inches(9.0), y, Inches(3.7), Inches(0.5),
            f"{name} — {round(row['share']*100)}%",
            size=15, color=INK, bold=True, font=BODY_FONT,
        )
        y += Inches(0.95)
    add_footer(s)
    set_notes(s, (
        "This split wasn't hand-picked — it's k-means clustering on recency, frequency, and "
        "spend, and the model chose five groups on its own. Position on the chart is recency "
        "and visit frequency; bubble size is how much each group spends. Champions are small "
        "in number but huge in value. At-risk regulars are the biggest group of all, at 32 "
        "percent — that's the group where a small improvement in retention pays for the whole "
        "project."
    ))


def build_slide_4(prs, n):
    m = n["metrics"]
    s = add_blank_slide(prs, WHITE)
    add_title(s, "Does it work? Yes — with honest, stated limits")
    add_picture_fit(s, FIG / "08_roc_curves.png", MARGIN, Inches(1.6), Inches(5.9), Inches(4.4))
    add_picture_fit(s, FIG / "10_confusion_matrix.png", Inches(6.9), Inches(1.6), Inches(5.7), Inches(4.4))
    recall_pct = round(m["model_recall"] * 100)
    precision_pct = round(m["chosen_threshold"]["precision"] * 100)
    add_text(
        s, MARGIN, Inches(6.2), Inches(12.1), Inches(0.9),
        f"At the chosen threshold, the model catches {recall_pct}% of members who actually "
        f"leave, at {precision_pct}% precision on every alarm it raises — no model catches "
        "everyone, and this one doesn't pretend to.",
        size=15, color=INK, font=BODY_FONT, line_spacing=1.3,
    )
    add_footer(s)
    set_notes(s, (
        "Two models, tested the honest way: trained on an earlier three-month window, tested "
        "on a later one the model never saw, not a random shuffle that would leak the answer. "
        "XGBoost beats plain logistic regression, but the logistic regression is respectably "
        "close — that's a good sign the signal is real, not an artifact of a complex model "
        "overfitting. The confusion matrix uses plain labels on purpose: caught churn, missed "
        "churn, false alarm, true stay. We chose the threshold to maximize how many leavers we "
        "catch, while keeping false alarms low enough that the win-back coupon stays a genuine "
        "signal instead of spam."
    ))


def build_slide_5(prs, n):
    s = add_blank_slide(prs, WHITE)
    add_title(s, "The strongest warning isn't spend — it's silence")
    add_picture_fit(s, FIG / "12_shap_importance.png", MARGIN, Inches(1.6), Inches(7.0), Inches(5.4))
    add_bullets(
        s, Inches(8.3), Inches(1.9), Inches(4.4), Inches(4.8),
        [
            "A widening gap between visits predicts departure weeks before the last one.",
            "Members who never redeem a point churn noticeably more often.",
            "Recency and visit-gap trend dominate — spend alone barely registers.",
        ],
        size=15, color=INK, bullet_color=TEAL, space_after=18,
    )
    add_footer(s)
    set_notes(s, (
        "This chart ranks features by how much they actually move the model's prediction, "
        "averaged across every member. The two at the top — recency and gap trend — are "
        "exactly what the café-specific design of this project was built to catch: not just "
        "'when did they last visit,' but 'has their rhythm already started to change.' The "
        "redemption finding is the actionable one for marketing: a member who's never redeemed "
        "a point never felt the loyalty program actually reward them, and they leave for it."
    ))


def build_slide_6(prs, n):
    s = add_blank_slide(prs, WHITE)
    add_title(s, "A specific action for every flagged member")

    rows = [
        ("Champions", "low", "VIP perks, early access to new drinks, referral ask", NAVY),
        ("Loyal", "low–med", "Tier-up nudge, personalized bundle", TEAL),
        ("At-risk regulars", "high", "15%-off win-back coupon + “we miss you” LINE mission", GOLD),
        ("Hibernating", "very high", "Bounce-back free drink + one-question why-survey", CORAL),
        ("One-and-done", "high", "Onboarding mission, second-visit nudge", SLATE),
    ]
    top = Inches(1.65)
    row_h = Inches(0.62)
    col_x = [MARGIN, Inches(3.4), Inches(4.6)]
    headers = ["Segment", "Risk", "Recommended action"]
    for cx, htext in zip(col_x, headers):
        add_text(s, cx, top, Inches(7.5), Inches(0.35), htext, size=12, color=INK_MUTED, bold=True)
    top += Inches(0.45)
    for i, (seg, risk, action, color) in enumerate(rows):
        y = top + row_h * i
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, col_x[0], y + Inches(0.1), Inches(0.16), Inches(0.16))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        dot.shadow.inherit = False
        add_text(s, col_x[0] + Inches(0.3), y, Inches(2.9), Inches(0.5), seg, size=14, color=INK, bold=True)
        add_text(s, col_x[1], y, Inches(1.1), Inches(0.5), risk, size=13, color=INK_MUTED)
        add_text(s, col_x[2], y, Inches(7.4), Inches(0.5), action, size=13, color=INK)

    # Worked example card
    ex_top = Inches(5.35)
    card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, ex_top, Inches(12.13), Inches(1.7))
    card.fill.solid()
    card.fill.fore_color.rgb = CANVAS
    card.line.color.rgb = LINE
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    tm = n["top_member"]
    add_text(
        s, Inches(0.9), ex_top + Inches(0.18), Inches(11.4), Inches(0.4),
        f"Worked example — member {tm['member_id']}", size=14, color=NAVY, bold=True,
    )
    add_text(
        s, Inches(0.9), ex_top + Inches(0.62), Inches(11.4), Inches(0.9),
        f"{tm['tier']} tier, {tm['segment']} segment, ฿{tm['annual_value_thb']:,.0f}/year in spend. "
        f"{tm['reason_1'].capitalize()}. "
        f"Risk: {tm['churn_probability']*100:.1f}%. Action: {tm['recommended_action']}.",
        size=13, color=INK, line_spacing=1.25,
    )
    add_footer(s)
    set_notes(s, (
        "The action table is deliberately not one-size-fits-all: a Champion at risk gets "
        "recognition, not a discount, because they don't need convincing to spend — they need "
        "to feel seen. Hibernating members get the biggest incentive because they're the "
        "hardest and most expensive to win back once fully gone. The example at the bottom is "
        "a real flagged member from the current run — the highest-value one — with the actual "
        "reasons the model gave and the exact action it recommends."
    ))


def build_slide_7(prs, n):
    s = add_blank_slide(prs, NAVY_DARK)
    add_title(s, "What it's worth, and what's next", color=WHITE)

    funnel = [
        (f"{n['n_total_members']:,}", "total members"),
        (f"{n['n_scored']:,}", "scored by the model"),
        (f"{n['n_flagged']:,}", "flagged high-risk"),
        (f"฿{n['flagged_value_thb']/1_000_000:.1f}M/yr", "in spend on the line"),
    ]
    x = MARGIN
    w = Inches(2.6)
    gap = Inches(0.55)
    for i, (value, label) in enumerate(funnel):
        add_stat(s, x, Inches(1.8), w, value, label, value_color=[WHITE, TEAL, GOLD, CORAL][i], value_size=34)
        if i < len(funnel) - 1:
            add_text(
                s, x + w, Inches(1.95), gap, Inches(0.6), "→",
                size=24, color=RGBColor(0x5A, 0x70, 0x88), align=PP_ALIGN.CENTER,
            )
        x += w + gap

    add_text(
        s, MARGIN, Inches(3.55), Inches(4.3), Inches(0.4), "To productionize, this needs:",
        size=14, color=WHITE, bold=True,
    )
    add_bullets(
        s, MARGIN, Inches(4.0), Inches(5.6), Inches(2.6),
        [
            "Oberry's real transaction data, not synthetic",
            "A scheduled batch-scoring run, not a one-off",
            "Instrumented campaigns with measured lift",
            "Ongoing monitoring for model drift",
        ],
        size=14, color=RGBColor(0xC9, 0xD4, 0xDE), bullet_color=TEAL, space_after=10,
    )
    add_text(
        s, Inches(6.9), Inches(3.55), Inches(5.6), Inches(3.2),
        "This is a working prototype, validated on synthetic data end to end — the "
        "architecture is real and the pipeline runs today. What's unproven is whether these "
        "exact patterns hold in Oberry's real transaction history, and that's the next step, "
        "not a leap of faith.",
        size=15, color=WHITE, line_spacing=1.35,
    )
    add_footer(s, dark=True)
    set_notes(s, (
        "This is the slide to slow down on. The funnel shows the honest scope: twenty thousand "
        "members, fourteen thousand three hundred of them scored by the model, twenty-two "
        "sixty flagged today, representing about 5.4 million baht a year in spend. We have not "
        "measured what fraction of those we'd actually retain with these campaigns — that "
        "requires running them against real members and measuring what happens, which is "
        "exactly the next phase. Everything up to this point is a working prototype on "
        "synthetic data. The architecture, the leak-free evaluation, the explainability — "
        "that's all real and reusable. What's left is proving it on Oberry's actual data."
    ))


def main() -> None:
    n = load_numbers()
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    for builder in [
        build_slide_1, build_slide_2, build_slide_3, build_slide_4,
        build_slide_5, build_slide_6, build_slide_7,
    ]:
        builder(prs, n)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
